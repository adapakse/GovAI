# -*- coding: utf-8 -*-
"""Generator prezentacji zapasowej GovAI — scenariusz demo TrustBank (11 aktów).

Zrzuty ekranu pochodzą z rzeczywistego środowiska demo (16.07.2026), nie są
mockupami. Używać wyłącznie jako plan B, gdyby live demo nie zadziałało —
patrz presentations/DEMO_SCENARIUSZ.md, sekcja E (plan awaryjny).
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta — zgodna z brand guide GovAI (branding/brand-guide.html) ───────────
DARK   = RGBColor(0x0D, 0x1B, 0x2A)
NAVY   = RGBColor(0x1B, 0x35, 0x5E)
BLUE   = RGBColor(0x1E, 0x6F, 0xBF)
TEAL   = RGBColor(0x00, 0xB4, 0xD8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY  = RGBColor(0xCB, 0xD5, 0xE1)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x3D, 0xC9, 0x7A)

LOGO_LOCKUP = r"C:\Users\Adam\govai\branding\raster\logo-lockup-dark-bg.png"
LOGO_MARK   = r"C:\Users\Adam\govai\branding\raster\logo-mark.png"
SHOTS       = r"C:\Users\Adam\Pictures\Screenshots\GovAI_screens"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=DARK):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s


def textbox(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
            align=PP_ALIGN.LEFT, font="Calibri", italic=False, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return box


def bullets(slide, l, t, w, h, items, size=14.5, color=LGRAY, space_after=10):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return box


def rect(slide, l, t, w, h, fill=BLUE, line=None, radius=False):
    from pptx.enum.shapes import MSO_SHAPE
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def kicker(slide, text, color=TEAL):
    textbox(slide, 0.6, 0.32, 10, 0.4, text.upper(), size=13, color=color, bold=True)


def footer(slide, n):
    slide.shapes.add_picture(LOGO_MARK, Inches(0.55), Inches(7.08), height=Inches(0.28))
    textbox(slide, 12.3, 7.07, 0.8, 0.35, str(n), size=11, color=MGRAY)
    rect(slide, 0, 7.02, 13.333, 0.02, fill=BLUE)


def title_bar(slide, title):
    textbox(slide, 0.6, 0.68, 12.1, 0.75, title, size=25, color=WHITE, bold=True)


def screenshot(slide, path, box_l=0.6, box_t=1.75, box_w=12.13, box_h=3.55, border=True):
    """Wstawia zrzut ekranu wpasowany (contain) w zadany prostokąt, wyśrodkowany."""
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    cand_w = box_h * aspect
    if cand_w <= box_w:
        disp_h, disp_w = box_h, cand_w
    else:
        disp_w, disp_h = box_w, box_w / aspect
    left = box_l + (box_w - disp_w) / 2
    top = box_t + (box_h - disp_h) / 2
    if border:
        pad = 0.04
        rect(slide, left - pad, top - pad, disp_w + 2 * pad, disp_h + 2 * pad,
             fill=NAVY, line=BLUE)
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))
    return top + disp_h  # dolna krawędź obrazu, do pozycjonowania kolejnych elementów


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def act_slide(n, act_label, title, img_file, bullet_items, note_text, act_color=TEAL,
              img_box=(0.6, 1.7, 12.13, 3.5)):
    s = add_slide()
    kicker(s, act_label, color=act_color)
    title_bar(s, title)
    img_path = os.path.join(SHOTS, img_file)
    bottom = screenshot(s, img_path, *img_box)
    bullets(s, 0.6, 5.55, 12.13, 1.4, bullet_items, size=13.5, space_after=6)
    footer(s, n)
    notes(s, note_text)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLAJD 1 — TYTUŁ
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 4.75, 13.333, 0.05, fill=TEAL)
textbox(s, 0.9, 1.75, 11, 0.5, "WERSJA ZAPASOWA DEMONSTRACJI", size=15, color=TEAL, bold=True)
s.shapes.add_picture(LOGO_LOCKUP, Inches(0.78), Inches(2.25), height=Inches(1.15))
textbox(s, 0.9, 3.65, 11.3, 1.0,
        "Scenariusz „TrustBank” — 3 agentów AI, pełny cykl bramki GovAI:\n"
        "rejestr → PII → polityki → nadzór człowieka → dziennik → raport zgodności",
        size=17, color=LGRAY, line_spacing=1.3)
textbox(s, 0.9, 5.0, 11, 0.5,
        "Zrzuty z rzeczywistego, działającego środowiska (16.07.2026) — do użycia, gdyby coś nie zadziałało na żywo.",
        size=13, color=MGRAY, italic=True)
notes(s, "Ta prezentacja to plan B dla scenariusza demo z DEMO_SCENARIUSZ.md. "
         "Używaj tylko jeśli środowisko live przestanie działać w trakcie spotkania. "
         "Każdy slajd odpowiada jednemu aktowi z tego samego scenariusza — możesz mówić "
         "dokładnie to, co mówiłbyś, pokazując aplikację na żywo.")

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 2 — NOTATKA DLA PROWADZĄCEGO
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Notatka dla prowadzącego")
title_bar(s, "Jak korzystać z tej prezentacji")
bullets(s, 0.6, 1.9, 12.1, 4.3, [
    "To jest wersja zapasowa — używaj tylko, jeśli aplikacja / internet / model AI zawiedzie na żywo.",
    "11 aktów, ten sam scenariusz co live demo: bank „TrustBank”, 3 agentów — Asystent Obsługi Klienta "
    "(ograniczone ryzyko), Agent Oceny Kredytowej (wysokie ryzyko), Agent Rekrutacji Wewnętrznej (wysokie ryzyko).",
    "Zrzuty pochodzą z rzeczywistego, działającego środowiska GovAI — to nie są mockupy ani makiety.",
    "Notatki prowadzącego (widok „Notatki” w PowerPoincie) zawierają pełny tekst do powiedzenia przy każdym akcie.",
    "Dane logowania (gdyby jednak udało się odpalić środowisko w trakcie): "
    "partner — admin@kancelaria.local, it_admin — it@kancelaria.local.",
], size=15.5, space_after=16)
footer(s, 2)
notes(s, "Krótkie wprowadzenie — nie czytać na głos widowni, to instrukcja tylko dla Ciebie.")

# ════════════════════════════════════════════════════════════════════════════
# AKT 1 — PULPIT
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    3, "Akt 1 · 2 min", "Kontrola z lotu ptaka",
    "pulpit 2026-07-16 074342.png",
    [
        "Pulpit agreguje ostatnie 7 dni z dziennika: 57 aktywnych agentów (20 wysokiego ryzyka), "
        "94 wywołania, 16 blokad (17%), 9 błędów infrastruktury, 0 oczekujących w nadzorze.",
        "Wykres „Wywołania wg agenta” pokazuje proporcję dozwolonych i zablokowanych wywołań na agenta.",
        "Panel „Zdarzenia real-time” — połączenie WebSocket, zdarzenia pojawiają się na żywo bez odświeżania strony.",
    ],
    "Każde wywołanie każdego agenta jest tu widoczne. X blokad, Y zamaskowanych danych osobowych, "
    "Z decyzji skierowanych do człowieka — wszystko bez dotykania kodu agentów.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 2 — REJESTR I KLASYFIKACJA AI ACT
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    4, "Akt 2 · 2 min", "System sam rozumie ryzyko regulacyjne",
    "Agent oceny kredytowej 2026-07-16 074526.png",
    [
        "Agent Oceny Kredytowej — poziom ryzyka WYSOKIE (Aneks III pkt 5b), status ACTIVE, oznaczony jako wymagający nadzoru.",
        "Podstawa prawna klasyfikacji przypisana automatycznie: „ocena zdolności kredytowej osób fizycznych”.",
        "Zakładka „Zgodność AI Act”: 8/9 wymagań spełnionych, 1 ważna luka do uzupełnienia — z artykułami "
        "(Art. 9, Art. 10…), opisem i statusem każdego wymogu.",
    ],
    "System sam sklasyfikował agenta jako wysokiego ryzyka i wie, jakie obowiązki z tego wynikają — "
    "to nie jest statyczny opis, tylko żywa ocena zgodności liczona z danych rejestru.",
)

# ════════════════════════════════════════════════════════════════════════════
# DIVIDER — SYMULATOR
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    5, "Symulator agentów", "Kolejne cztery akty dzieją się w jednym miejscu",
    "Symulator 2026-07-16 074631.png",
    [
        "Każde wywołanie przechodzi przez ten sam pipeline co produkcyjny agent: GovAI Gateway → skan PII → "
        "silnik polityk → model AI.",
        "Gotowe scenariusze przypisane do każdego z 3 agentów demo — uruchomienie jednym kliknięciem „Uruchom”.",
        "Wynik widoczny natychmiast tu, w Dzienniku audytowym i w kolejce Nadzoru.",
    ],
    "To jest bezpieczny „poligon” — pokazuje dokładnie to samo zachowanie bramki, co realne wywołanie agenta, "
    "bez potrzeby integrowania żadnego zewnętrznego systemu.",
    act_color=BLUE,
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 3 — WYWOŁANIE CZYSTE (ALLOWED)
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    6, "Akt 3 · 1–2 min", "Happy path — bramka nie przeszkadza, gdy wszystko jest OK",
    "Agent obsługi Zapytanie o ubezpieczenie 2026-07-16 074750.png",
    [
        "Scenariusz „Zapytanie o ubezpieczenie” na Asystencie Obsługi Klienta.",
        "Brak PII, brak naruszeń polityk → wybór providera wg wrażliwości danych → odpowiedź modelu (claude-haiku-4-5).",
        "HTTP 200, status zgodny z oczekiwanym — pełna treść odpowiedzi widoczna od razu w karcie scenariusza.",
    ],
    "To jest happy path: gdy zapytanie jest czyste, bramka w ogóle nie przeszkadza — tylko obserwuje i loguje.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 4 — OCHRONA PII
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    7, "Akt 4 · 2 min", "RODO/PII w praktyce — model nigdy nie widzi surowego PESEL-u",
    "Agent obsługi Zapytanie o ubezpieczenie 2026-07-16 074750.png",
    [
        "Kolejny scenariusz na tym samym agencie: „Klient podaje PESEL i numer konta (PII)” — oczekiwany wynik: "
        "Dozwolone + PII (widoczny poniżej pierwszego, jeszcze przed uruchomieniem).",
        "Presidio wykrywa PESEL i IBAN i maskuje je, zanim cokolwiek trafi do modelu — model widzi „[PESEL_REDACTED]”.",
        "Wpis w Dzienniku odnotowuje wykryte kategorie PII przy zachowanym statusie „Dozwolone”.",
    ],
    "Model nigdy nie zobaczył surowego PESEL-u — został zamaskowany w bramce, zanim opuścił naszą infrastrukturę. "
    "Na żywo: kliknij „Uruchom” przy tym scenariuszu, żeby pokazać zamaskowaną odpowiedź.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 5 — BLOKADA G-001
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    8, "Akt 5 · 1–2 min", "Polityka bezpieczeństwa zatrzymuje niebezpieczną operację",
    "Próba mutacji finansowej 2026-07-16 074852.png",
    [
        "Scenariusz „Próba mutacji finansowej” — fraza „zmień saldo” dopasowuje regułę G-001.",
        "Blokada następuje PRZED wywołaniem modelu — model w ogóle nie widzi tego zapytania.",
        "HTTP 403, powód zapisany wprost: „Modyfikacja danych finansowych poza zakresem agenta”.",
    ],
    "Klasyczne nadużycie — agent obsługi klienta nie ma prawa modyfikować sald. Polityka złapała to, "
    "zanim zapytanie dotarło do modelu.",
    act_color=RED,
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 6 — BLOKADA G-002 (PROMPT INJECTION)
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    9, "Akt 6 · 1 min", "Klasyczny atak na agenta — zatrzymany, zanim dotarł do modelu",
    "Promt injection 2026-07-16 074946.png",
    [
        "Scenariusz „Atak prompt injection” — fraza „ignore previous instructions” dopasowuje regułę G-002.",
        "HTTP 403, natychmiastowa blokada: „Wykryto próbę wstrzyknięcia instrukcji (prompt injection)”.",
        "Ta sama reguła chroni też agenta rekrutacji — scenariusz „CV z atakiem injection” działa identycznie.",
    ],
    "Prompt injection to jeden z najczęstszych ataków na agentów AI. Reguła G-002 łapie go na poziomie bramki, "
    "niezależnie od tego, który agent jest celem.",
    act_color=RED,
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 7a — NADZÓR CZŁOWIEKA: ZGŁOSZENIE
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    10, "Akt 7 · punkt kulminacyjny", "Agent wysokiego ryzyka nie kończy decyzji sam",
    "Nadzór kolejka 2026-07-16 075216.png",
    [
        "Scenariusz „Dobry wniosek kredytowy” na Agencie Oceny Kredytowej → decyzja skierowana do nadzoru "
        "(Oversight ID w odpowiedzi gatewaya, widoczny wcześniej w Symulatorze).",
        "Kolejka Nadzór pokazuje pozycję oczekującą z licznikiem czasu (tu: 58m 38s) — eskalacja, jeśli nikt nie zatwierdzi.",
        "Baner ostrzegawczy: czas przeglądu każdego zadania jest mierzony od momentu otwarcia karty (art. 14 EU AI Act).",
    ],
    "Agent wysokiego ryzyka nie kończy decyzji sam — wywołanie trafia do kolejki nadzoru z licznikiem TTL. "
    "Otwórz pozycję i kliknij „Przejrzyj”.",
    act_color=ORANGE,
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 7b — NADZÓR CZŁOWIEKA: DECYZJA
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    11, "Akt 7 (cd.)", "Nie tylko wymusza nadzór — pilnuje, czy jest realny",
    "Nadzór przejrzyj 2026-07-16 075245.png",
    [
        "Recenzent widzi pełną treść wniosku, typ decyzji („Wywołanie agenta HIGH”) i czas zgłoszenia.",
        "Trzy opcje decyzji: Zatwierdź / Odrzuć / Eskaluj — plus pole uzasadnienia, które jest wymagane.",
        "Ostrzeżenie wprost w interfejsie: zatwierdzenie poniżej 10 sekund od otwarcia karty jest oznaczane "
        "jako „pozorny nadzór” i zapisywane w dzienniku audytowym.",
    ],
    "System nie tylko wymusza nadzór — pilnuje też, czy nadzór jest realny. Jeśli zatwierdzisz zbyt szybko, "
    "zobaczysz alert rubber-stamp. Zatwierdź normalnie, żeby pokazać ślad w Dzienniku.",
    act_color=ORANGE,
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 8 — POLITYKI BEZ KODU
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    12, "Akt 8 · 2 min", "Zgodność konfigurowalna przez biznes, nie przez programistów",
    "Polityki lista słów kluczowych 2026-07-16 075330.png",
    [
        "Reguły G-001 (mutacja finansowa) i G-002 (prompt injection) — lista słów kluczowych edytowalna wprost w panelu.",
        "Gateway odświeża reguły z bazy co 60 sekund automatycznie — zero wdrożenia, zero restartu.",
        "Priorytet i zasięg (org / agent) decydują o kolejności sprawdzania reguł blokujących.",
    ],
    "Compliance officer może dodać nowe słowo kluczowe (np. „przelej wszystkie środki”) i zapisać — "
    "po chwili nowa fraza zaczyna blokować. Bez wdrożenia, bez programisty.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 9 — ROUTING WIELODOSTAWCOWY
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    13, "Akt 9 · 1–2 min", "Dane wrażliwe nie wychodzą tam, gdzie nie powinny",
    "Prowiderzy 2026-07-16 075416.png",
    [
        "5 providerów, wszystkie 5 zdrowe — w tym „Bielik On-Prem (Local)” jako jedyny provider poziomu "
        "„Uprzywilejowany”, oparty o lokalny model Ollama.",
        "Gateway klasyfikuje wrażliwość zapytania (public → internal → confidential → privileged) i wybiera "
        "aktywnego providera o najniższym priorytecie, który obsługuje wymagany poziom.",
        "Dane „privileged” (np. tajemnica adwokacka) trafiają wyłącznie do providerów on-prem — nigdy do chmury publicznej.",
    ],
    "Ten sam agent, ale dane poufne trafią do innego, bezpieczniejszego modelu — automatycznie, bez zmiany "
    "kodu agenta.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 10 — PARAMETRY BEZ WDROŻENIA
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    14, "Akt 10 · 1–2 min · nowość", "Strojenie systemu w locie przez administratora",
    "PArametry 2026-07-16 075513.png",
    [
        "Progi, stawki cennika (USD/1k tokenów) i budżety domyślne trzymane w bazie — zmiana propaguje do "
        "API i bramki w ≤60 sekund.",
        "RBAC egzekwowany w interfejsie: edycja wymaga roli it_admin. Partner widzi panel w „trybie podglądu” "
        "(widać, ale nie da się edytować).",
        "Każdy parametr ma historię zmian (kto, kiedy, stara/nowa wartość) — link „historia” przy każdej pozycji.",
    ],
    "Tylko IT może zmieniać parametry techniczne — role są egzekwowane na poziomie interfejsu i API. "
    "Uwaga: na tym zrzucie widać kosmetyczny błąd kodowania polskich znaków w opisach pól — do naprawy "
    "przed spotkaniem, nie wpływa na działanie systemu.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 11a — DZIENNIK AUDYTOWY
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    15, "Akt 11 · domknięcie", "Niezmienny ślad każdego wywołania",
    "Dziennik audytowy 2026-07-16 075552.png",
    [
        "98 wpisów, filtrowalne wg wyniku, wykrycia PII i zakresu czasu.",
        "Każdy akt z tego demo zostawił tu wpis: OK, BLOKADA, NADZÓR — z agentem, typem zdarzenia i latencją.",
        "Status „BŁĄD” (fioletowy) — nowość: pokrywa awarie infrastrukturalne (brak providera, błąd modelu), "
        "które wcześniej ginęły bez śladu w dzienniku.",
    ],
    "Pełna ścieżka audytowa — każdy akt, który właśnie pokazaliśmy, zostawił tu niezmienny wpis.",
)

# ════════════════════════════════════════════════════════════════════════════
# AKT 11b — RAPORT ZGODNOŚCI
# ════════════════════════════════════════════════════════════════════════════
act_slide(
    16, "Akt 11 (cd.)", "Gotowy raport dla regulatora — jeden klik",
    "Agent Raport 2026-07-16 075752.png",
    [
        "Raport AI Act per agent: poziom ryzyka, kategoria Aneksu III, wymóg nadzoru, streszczenie wykonawcze "
        "wygenerowane automatycznie na podstawie danych operacyjnych.",
        "Statystyki 30 dni wprost w raporcie: liczba wywołań, blokad, decyzji skierowanych do nadzoru.",
        "Eksport do PDF jednym kliknięciem („Pobierz PDF”) — gotowy dokument do przekazania regulatorowi.",
    ],
    "Pełna ścieżka audytowa i gotowy raport dla regulatora — dowód, że firma panuje nad swoimi agentami.",
    act_color=GREEN,
)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 17 — PUENTA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
textbox(s, 0.9, 2.3, 11.5, 0.5, "PUENTA", size=15, color=TEAL, bold=True)
textbox(s, 0.9, 2.85, 11.5, 2.3,
        "„Pokazaliśmy jeden bank z trzema agentami. GovAI skaluje to na setki agentów:\n"
        "jedna bramka, jeden dziennik, jedna konsola — pełna zgodność z EU AI Act bez\n"
        "przepisywania kodu agentów. To różnica między *mamy AI* a *panujemy nad AI*.”",
        size=22, color=WHITE, bold=True, italic=True, line_spacing=1.3)
textbox(s, 0.9, 5.6, 11, 0.5,
        "Ta prezentacja jest kopią zapasową — pełne, interaktywne demo dostępne pod http://localhost:4000",
        size=13, color=MGRAY, italic=True)
s.shapes.add_picture(LOGO_LOCKUP, Inches(0.9), Inches(6.55), height=Inches(0.42))
footer(s, 17)
notes(s, "Zamknięcie — to samo, co na koniec live demo. Jeśli był to fallback z powodu awarii, "
         "warto dodać zdanie: „To była wersja zapasowa — środowisko na żywo pokażemy przy najbliższej okazji.”")

OUT = r"C:\Users\Adam\govai\presentations\GovAI_Demo_Backup.pptx"
prs.save(OUT)
print(f"OK — zapisano prezentację: {OUT}")
