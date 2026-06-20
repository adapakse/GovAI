# -*- coding: utf-8 -*-
"""Generator prezentacji GovAI — rynek: finanse, kancelarie prawne, medycyna."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paleta — zgodna z brand guide GovAI (branding/brand-guide.html) ───────────
DARK   = RGBColor(0x0D, 0x1B, 0x2A)
NAVY   = RGBColor(0x1B, 0x35, 0x5E)
BLUE   = RGBColor(0x1E, 0x6F, 0xBF)
TEAL   = RGBColor(0x00, 0xB4, 0xD8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY  = RGBColor(0xCB, 0xD5, 0xE1)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x3D, 0xC9, 0x7A)
YELLOW = RGBColor(0xE8, 0xC5, 0x47)

LOGO_LOCKUP = r"C:\Users\Adam\govai\branding\raster\logo-lockup-dark-bg.png"
LOGO_MARK   = r"C:\Users\Adam\govai\branding\raster\logo-mark.png"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=DARK):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s


def textbox(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
            align=PP_ALIGN.LEFT, font="Calibri", italic=False, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return box


def bullets(slide, l, t, w, h, items, size=15, color=LGRAY, bold_first=False,
            space_after=10, bullet_color=TEAL):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return box


def rect(slide, l, t, w, h, fill=BLUE, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def kicker(slide, text, color=TEAL):
    textbox(slide, 0.6, 0.35, 8, 0.4, text.upper(), size=13, color=color, bold=True)


def footer(slide, n):
    slide.shapes.add_picture(LOGO_MARK, Inches(0.55), Inches(7.06), height=Inches(0.3))
    textbox(slide, 12.3, 7.05, 0.8, 0.35, str(n), size=11, color=MGRAY)
    rect(slide, 0, 7.0, 13.333, 0.02, fill=BLUE)


def title_bar(slide, title, sub=None):
    textbox(slide, 0.6, 0.55, 12, 0.8, title, size=30, color=WHITE, bold=True)
    if sub:
        textbox(slide, 0.6, 1.15, 12, 0.5, sub, size=15, color=MGRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLAJD 1 — TYTUŁ
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 4.55, 13.333, 0.05, fill=TEAL)
textbox(s, 0.9, 2.05, 11, 0.5, "EU AI ACT  ·  ZGODNOŚĆ I KONTROLA SYSTEMÓW AI", size=15,
        color=TEAL, bold=True)
s.shapes.add_picture(LOGO_LOCKUP, Inches(0.78), Inches(2.55), height=Inches(1.25))
textbox(s, 0.9, 4.75, 11.3, 0.7,
        "Bramka kontrolna i rejestr zgodności dla agentów AI — finanse, rekrutacja, kancelarie prawne, ochrona zdrowia i inne sektory wysokiego ryzyka",
        size=18, color=LGRAY)
textbox(s, 0.9, 6.5, 11, 0.5, "Prezentacja rozwiązania · Partnerstwo wdrożeniowe", size=13, color=MGRAY, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 2 — KONTEKST: CO TO JEST AI ACT
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Kontekst regulacyjny")
title_bar(s, "EU AI Act — pierwsze na świecie horyzontalne prawo o sztucznej inteligencji")
bullets(s, 0.6, 1.9, 6.0, 4.6, [
    "Rozporządzenie UE 2024/1689 — weszło w życie 1 sierpnia 2024",
    "Obowiązuje bezpośrednio we wszystkich 27 państwach członkowskich — bez potrzeby ustaw krajowych",
    "Podejście oparte na ryzyku: minimalne / ograniczone / wysokie / niedopuszczalne",
    "Obowiązki dotyczą zarówno dostawców (providers), jak i użytkowników biznesowych (deployers) systemów AI",
    "Zasięg eksterytorialny — dotyczy każdej firmy, której system AI wpływa na osoby w UE, niezależnie od siedziby dostawcy modelu",
], size=15.5)
rect(s, 7.0, 1.85, 5.7, 4.7, fill=BLUE, radius=True)
textbox(s, 7.35, 2.05, 5.0, 0.4, "HARMONOGRAM WDROŻENIA", size=13, color=TEAL, bold=True)
timeline = [
    ("Luty 2025", "Zakaz praktyk niedopuszczalnego ryzyka (manipulacja, scoring społeczny, biometria masowa)"),
    ("Sierpień 2025", "Obowiązki dla modeli ogólnego przeznaczenia (GPAI) — Claude, GPT, DeepSeek i inne"),
    ("Sierpień 2026", "Pełne obowiązki dla systemów wysokiego ryzyka (Aneks III) — KONTROLE REGULATORA"),
    ("Sierpień 2027", "Systemy wysokiego ryzyka w produktach już regulowanych (np. wyroby medyczne)"),
]
y = 2.55
for date, desc in timeline:
    rect(s, 7.35, y, 0.09, 0.62, fill=TEAL)
    textbox(s, 7.6, y - 0.04, 1.7, 0.3, date, size=13.5, color=WHITE, bold=True)
    textbox(s, 7.6, y + 0.26, 4.5, 0.6, desc, size=11.5, color=LGRAY, line_spacing=1.05)
    y += 1.0
footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 3 — HARMONOGRAM WEJŚCIA W ŻYCIE I WDROŻENIE W POLSCE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Kontekst regulacyjny")
title_bar(s, "AI Act wchodzi w życie etapami — i już teraz staje się prawem polskim")
textbox(s, 0.6, 1.95, 12.1, 0.55,
        "Główne przepisy Aktu o sztucznej inteligencji zaczynają obowiązywać od 2 sierpnia 2026 roku — ale rozporządzenie wdrażane jest stopniowo:",
        size=14, color=LGRAY)
stages = [
    ("Sierpień 2025", "Weszły w życie regulacje dla dostawców modeli AI ogólnego przeznaczenia (GPAI), kary administracyjne oraz obowiązek zapewnienia kompetencji w zakresie AI w organizacji.", BLUE),
    ("2 sierpnia 2026", "Zaczyna obowiązywać większość przepisów i sankcji w pełnej skali — w tym przejrzystość i ocena ryzyka dla systemów wysokiego ryzyka.", ORANGE),
    ("2027", "Pełne wejście w życie pozostałych wymogów, w tym dla produktów już regulowanych sektorowo (np. wyroby medyczne).", TEAL),
]
y = 2.65
for date, desc, color in stages:
    rect(s, 0.6, y, 12.1, 1.05, fill=NAVY, radius=True)
    rect(s, 0.6, y, 0.12, 1.05, fill=color)
    textbox(s, 1.0, y + 0.12, 2.6, 0.8, date, size=15, color=color, bold=True)
    textbox(s, 3.75, y + 0.12, 8.7, 0.85, desc, size=12, color=LGRAY, line_spacing=1.15)
    y += 1.2
rect(s, 0.6, y + 0.1, 12.1, 1.15, fill=RGBColor(0x10, 0x2E, 0x2A), radius=True, line=GREEN)
textbox(s, 0.95, y + 0.25, 11.4, 0.4, "POLSKA — STAN PRAC LEGISLACYJNYCH", size=12.5, color=GREEN, bold=True)
textbox(s, 0.95, y + 0.6, 11.4, 0.55,
        "Trwają ostateczne prace nad ustawą krajową wdrażającą AI Act — projekt trafił pod obrady i został uchwalony przez Sejm w czerwcu 2026 roku.",
        size=12, color=LGRAY, line_spacing=1.15)
footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 4 — KARY I SANKCJE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Ryzyko finansowe", color=RED)
title_bar(s, "Kary administracyjne — wśród najwyższych w prawie unijnym")
cards = [
    ("do 35 mln EUR", "lub 7% globalnego obrotu rocznego", "Praktyki niedopuszczalnego ryzyka (art. 5) — np. nielegalny scoring społeczny, manipulacyjne systemy AI", RED),
    ("do 15 mln EUR", "lub 3% globalnego obrotu rocznego", "Naruszenie obowiązków dla systemów wysokiego ryzyka — dokumentacja, nadzór człowieka, ocena zgodności", ORANGE),
    ("do 7,5 mln EUR", "lub 1% globalnego obrotu rocznego", "Wprowadzenie organów nadzoru w błąd — niepełne, nieprawidłowe informacje przy kontroli", YELLOW),
]
x = 0.6
for amount, basis, desc, color in cards:
    rect(s, x, 1.95, 4.0, 4.35, fill=BLUE, radius=True)
    rect(s, x, 1.95, 4.0, 0.12, fill=color, radius=False)
    textbox(s, x + 0.3, 2.3, 3.4, 0.7, amount, size=26, color=WHITE, bold=True)
    textbox(s, x + 0.3, 2.95, 3.4, 0.5, basis, size=12.5, color=color, bold=True)
    textbox(s, x + 0.3, 3.6, 3.4, 2.5, desc, size=13, color=LGRAY, line_spacing=1.2)
    x += 4.25
textbox(s, 0.6, 6.45, 12, 0.5,
        "Stosuje się WYŻSZĄ z dwóch wartości — dla globalnych grup finansowych i sieci kancelarii liczy się obrót całej grupy, nie tylko spółki lokalnej.",
        size=12.5, color=MGRAY, italic=True)
footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 5 — KONSEKWENCJE POZA KARAMI
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Poza karami finansowymi", color=ORANGE)
title_bar(s, "Konsekwencje, które nie mieszczą się w arkuszu kalkulacyjnym")
items = [
    ("⛔", "Wstrzymanie systemu", "Regulator może nakazać natychmiastowe wycofanie systemu AI z rynku lub zaprzestanie jego używania — utrata zdolności operacyjnej z dnia na dzień."),
    ("⚖", "Odpowiedzialność cywilna", "Klient poszkodowany błędną decyzją AI (odmowa kredytu, błędna porada prawna, błędna diagnoza) może dochodzić odszkodowania — AI Act ułatwia dowodzenie winy dzięki obowiązkowej dokumentacji."),
    ("📉", "Ryzyko reputacyjne", "Sektor finansowy, prawny i medyczny opiera się na zaufaniu. Publiczna kontrola regulatora z negatywnym wynikiem trafia do mediów branżowych i prasy."),
    ("🔒", "Utrata licencji / akredytacji", "W sektorze medycznym i prawnym naruszenia AI Act mogą się nakładać na naruszenia regulacji branżowych (MDR, IVDR, przepisy o tajemnicy zawodowej) — podwójna ekspozycja."),
    ("⏳", "Zatrzymanie wdrożeń AI", "Firma bez udokumentowanej zgodności nie może bezpiecznie skalować użycia AI — każdy nowy projekt wymaga oceny ryzyka od zera, co kosztuje czas i pieniądze."),
]
y = 1.95
for icon, title, desc in items:
    rect(s, 0.6, y, 0.7, 0.7, fill=BLUE, radius=True)
    textbox(s, 0.6, y + 0.08, 0.7, 0.6, icon, size=24, color=TEAL, align=PP_ALIGN.CENTER)
    textbox(s, 1.5, y - 0.02, 2.8, 0.5, title, size=15, color=WHITE, bold=True)
    textbox(s, 4.4, y - 0.02, 8.3, 0.85, desc, size=12.5, color=LGRAY, line_spacing=1.1)
    y += 0.95
footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 6 — CZTERY PRZYKŁADOWE BRANŻE: PRZEGLĄD
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Rynek docelowy")
title_bar(s, "Cztery przykładowe branże, jedno wspólne ryzyko regulacyjne")
sectors = [
    ("🏦", "FINANSE", "Banki, ubezpieczyciele, fintech", [
        "Ocena zdolności kredytowej (Aneks III pkt 5b)",
        "Wykrywanie fraudów i blokady kont",
        "Scoring ryzyka ubezpieczeniowego",
        "Doradztwo inwestycyjne wspierane AI",
    ], BLUE),
    ("🧑‍💼", "REKRUTACJA", "Dział HR, agencje rekrutacyjne, RPO", [
        "Selekcja i ranking kandydatów (Aneks III pkt 4a)",
        "Automatyczne odrzucanie aplikacji bez przeglądu człowieka",
        "Analiza rozmów rekrutacyjnych (AI interview scoring)",
        "Ryzyko dyskryminacji = AI Act + prawo pracy naraz",
    ], BLUE),
    ("⚖", "KANCELARIE PRAWNE", "Doradztwo prawne, compliance, due diligence", [
        "Wsparcie interpretacji prawa i orzecznictwa (Aneks III pkt 8)",
        "Analiza umów i ryzyk kontraktowych AI",
        "Due diligence przy fuzjach i przejęciach",
        "Tajemnica zawodowa + RODO + AI Act = ekspozycja podwójna",
    ], BLUE),
    ("🏥", "MEDYCYNA", "Szpitale, diagnostyka, telemedycyna", [
        "Systemy wspierające diagnozę (Aneks III + MDR/IVDR)",
        "Triage i priorytetyzacja pacjentów",
        "Wyroby medyczne z komponentem AI — automatycznie wysokie ryzyko",
        "Błąd systemu = zagrożenie życia + odpowiedzialność karna",
    ], BLUE),
]
x = 0.6
w = 2.85
for icon, name, sub, points, color in sectors:
    rect(s, x, 1.9, w, 4.85, fill=color, radius=True)
    textbox(s, x + 0.2, 2.05, w - 0.4, 0.6, icon, size=26, color=TEAL)
    textbox(s, x + 0.2, 2.6, w - 0.4, 0.65, name, size=14.5, color=WHITE, bold=True)
    textbox(s, x + 0.2, 3.15, w - 0.4, 0.55, sub, size=10, color=MGRAY, italic=True)
    bullets(s, x + 0.2, 3.7, w - 0.35, 3.0, points, size=9.5, space_after=7)
    x += w + 0.2
footer(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 7 — DLACZEGO TE BRANŻE SĄ NAJBARDZIEJ NARAŻONE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Dlaczego teraz", color=ORANGE)
title_bar(s, "Wspólny mianownik: decyzje AI dotykają praw podstawowych obywateli")
textbox(s, 0.6, 1.85, 12, 0.55,
        "Aneks III AI Act wprost klasyfikuje jako „wysokie ryzyko” systemy używane w obszarach o silnym wpływie na jednostkę:",
        size=14, color=LGRAY)
table_rows = [
    ("Dostęp do usług finansowych", "Ocena kredytowa, scoring ubezpieczeniowy — bezpośredni wpływ na zdolność obywatela do uzyskania kredytu czy ubezpieczenia."),
    ("Zatrudnienie i selekcja kandydatów", "Automatyczne odrzucanie CV, ranking kandydatów, ocena wydajności pracowników — wpływa na dostęp do pracy i równe traktowanie (Aneks III pkt 4a)."),
    ("Wymiar sprawiedliwości i prawo", "Systemy wspierające interpretację i stosowanie prawa — błąd przenosi się na decyzje sądowe lub porady prawne klientów."),
    ("Ochrona zdrowia i życia", "Systemy w wyrobach medycznych i diagnostyce — błąd może prowadzić do uszczerbku na zdrowiu lub śmierci."),
]
y = 2.5
for title, desc in table_rows:
    rect(s, 0.6, y, 12.1, 0.95, fill=BLUE, radius=True)
    textbox(s, 0.9, y + 0.1, 4.0, 0.75, title, size=14, color=TEAL, bold=True)
    textbox(s, 5.1, y + 0.1, 7.4, 0.8, desc, size=11.5, color=LGRAY, line_spacing=1.1)
    y += 1.08
footer(s, 7)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 8 — TO NIE JEST LISTA ZAMKNIĘTA: INNE BRANŻE W OBSZARZE RYZYKA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Szerszy obraz", color=ORANGE)
title_bar(s, "Finanse, rekrutacja, prawo i medycyna to przykłady — nie cała lista")
textbox(s, 0.6, 1.9, 12.1, 0.6,
        "Aneks III AI Act obejmuje znacznie szerszy zakres sektorów. Każda organizacja wdrażająca AI w obszarach wpływających na ludzi powinna ocenić swoje ryzyko:",
        size=13.5, color=LGRAY)
other_sectors = [
    ("🎓", "Edukacja", "Ocena egzaminów, kwalifikacja na studia, monitorowanie uczniów podczas testów"),
    ("⚡", "Energetyka", "Zarządzanie krytyczną infrastrukturą energetyczną i sieciami przesyłowymi"),
    ("🏭", "Przemysł", "Bezpieczeństwo maszyn i systemów produkcyjnych z komponentem AI"),
    ("🚆", "Transport i infrastruktura", "Zarządzanie ruchem, systemy bezpieczeństwa w transporcie krytycznym"),
    ("🏛", "Sektor publiczny", "Przyznawanie świadczeń socjalnych, ocena uprawnień obywateli"),
    ("🛂", "Migracja i granice", "Kontrola graniczna, ocena wniosków azylowych i wizowych"),
]
cols = 3
cw, ch, gx, gy = 3.9, 1.6, 0.2, 0.25
x0, y0 = 0.6, 2.75
for i, (icon, name, desc) in enumerate(other_sectors):
    col, row = i % cols, i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=NAVY, radius=True)
    textbox(s, x + 0.2, y + 0.12, 0.7, 0.5, icon, size=22, color=TEAL)
    textbox(s, x + 0.2, y + 0.62, cw - 0.4, 0.35, name, size=13, color=WHITE, bold=True)
    textbox(s, x + 0.2, y + 0.98, cw - 0.4, 0.55, desc, size=9.5, color=LGRAY, line_spacing=1.05)
textbox(s, 0.6, 6.55, 12.1, 0.45,
        "Wspólny mianownik: systemy AI wpływające na dostęp do zasobów, praw i bezpieczeństwa obywateli — niezależnie od branży.",
        size=11.5, color=TEAL, italic=True, bold=True)
footer(s, 8)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 9 — ROZWIĄZANIE: ARCHITEKTURA GOVAI
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Rozwiązanie")
title_bar(s, "GovAI — bramka kontrolna między aplikacją a modelem AI")
# pipeline diagram
steps = [
    ("1", "Rejestr Agentów", "Każdy agent AI ma właściciela, opis, poziom ryzyka, podstawę prawną"),
    ("2", "Skan PII", "Presidio wykrywa i maskuje dane osobowe przed wysłaniem do modelu"),
    ("3", "Polityki", "Reguły blokujące niedozwolone operacje — w czasie rzeczywistym"),
    ("4", "Model AI", "Claude, DeepSeek, modele lokalne (Bielik) — niezależnie od dostawcy"),
    ("5", "Nadzór człowieka", "Decyzje wysokiego ryzyka czekają na zatwierdzenie recenzenta"),
    ("6", "Dziennik audytowy", "Każde wywołanie — kto, kiedy, na jakiej podstawie, z jakim wynikiem"),
]
x = 0.6
w = 1.95
for i, (num, title, desc) in enumerate(steps):
    rect(s, x, 2.5, w, 3.6, fill=BLUE, radius=True)
    rect(s, x + 0.65, 2.7, 0.65, 0.65, fill=TEAL, radius=True)
    textbox(s, x + 0.65, 2.75, 0.65, 0.55, num, size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + 0.12, 3.55, w - 0.24, 0.7, title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + 0.12, 4.25, w - 0.24, 1.7, desc, size=10.5, color=LGRAY, align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < len(steps) - 1:
        textbox(s, x + w - 0.05, 4.0, 0.3, 0.5, "→", size=22, color=TEAL, bold=True)
    x += w + 0.1
textbox(s, 0.6, 6.4, 12, 0.5,
        "Architektura niezależna od dostawcy modelu — działa identycznie dla Claude, DeepSeek czy modeli on-premise.",
        size=12.5, color=MGRAY, italic=True)
footer(s, 9)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 10 — KONTROLA KOSZTÓW I ALOKACJA MODELI
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Rozwiązanie")
title_bar(s, "Zgodność nie musi oznaczać najwyższego rachunku za AI")
textbox(s, 0.6, 2.05, 12.1, 0.6,
        "GovAI kieruje każde zapytanie do modelu odpowiedniego do poziomu ryzyka i złożoności zadania — nie zawsze do najdroższego.",
        size=13, color=MGRAY)
cost_cards = [
    ("Routing modelem wg ryzyka", BLUE,
     "Agent o minimalnym ryzyku (np. podsumowanie dokumentu) trafia do tańszego modelu (DeepSeek, Bielik). Zadania wysokiego ryzyka — do modeli klasy Claude/GPT z pełnym nadzorem."),
    ("Budżety i alerty na poziomie agenta", TEAL,
     "Każdy agent ma własny miesięczny budżet (monthly_budget_eur) i próg ostrzegawczy (cost_alert_threshold_eur) — przekroczenie generuje alert, zanim koszty wymkną się spod kontroli."),
    ("Koszt widoczny w czasie rzeczywistym", ORANGE,
     "Bramka liczy koszt każdego wywołania (tokeny × cennik dostawcy) i agreguje go na pulpicie oraz w raportach — KPI „Koszt API” na poziomie firmy, działu i agenta."),
]
y = 2.85
for title, color, desc in cost_cards:
    rect(s, 0.6, y, 12.1, 1.3, fill=NAVY, radius=True)
    rect(s, 0.6, y, 0.12, 1.3, fill=color)
    textbox(s, 1.0, y + 0.12, 11.3, 0.4, title, size=14.5, color=WHITE, bold=True)
    textbox(s, 1.0, y + 0.52, 11.3, 0.7, desc, size=11.5, color=LGRAY, line_spacing=1.15)
    y += 1.45
textbox(s, 0.6, 6.55, 12.1, 0.5,
        "Efekt: pełna zgodność i audytowalność przy koszcie proporcjonalnym do rzeczywistego ryzyka — nie do liczby zapytań.",
        size=12, color=TEAL, italic=True, bold=True)
footer(s, 10)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 11 — DEMO INTRO
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
textbox(s, 0.9, 2.6, 11.5, 0.6, "DEMONSTRACJA NA ŻYWO", size=16, color=TEAL, bold=True)
textbox(s, 0.9, 3.15, 11.5, 1.8,
        "„Bank z 340 agentami AI dostaje za 6 miesięcy\nkontrolę regulatora. Co wie o swoich agentach,\nco przetwarzają i kto ponosi odpowiedzialność?”",
        size=24, color=WHITE, bold=True, line_spacing=1.25)
textbox(s, 0.9, 5.4, 10, 0.5, "Pokażemy to na żywo — w czterech aktach.", size=15, color=MGRAY, italic=True)
footer(s, 11)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 12 — AKT 1: CHAOS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Akt 1", color=RED)
title_bar(s, "Chaos bez GovAI")
bullets(s, 0.6, 2.0, 11.8, 3.5, [
    "Trzy agenty AI wywołują model bezpośrednio — bez żadnej bramki kontrolnej",
    "Agent kredytowy wysyła numer PESEL klienta prosto do zewnętrznego modelu AI",
    "Agent rekrutacyjny podejmuje decyzję o kandydacie bez żadnego śladu audytowego",
    "Nikt w organizacji nie wie: ile to kosztuje, ile razy było wywołane, jakie dane wyciekły",
], size=16.5, color=LGRAY, space_after=18)
rect(s, 0.6, 5.6, 12.1, 1.0, fill=RGBColor(0x3A, 0x1A, 0x1A), radius=True)
textbox(s, 0.9, 5.85, 11.5, 0.5, "To jest dokładnie stan, w jakim regulator zastaje większość firm dziś.", size=14, color=RED, bold=True, italic=True)
footer(s, 12)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 13 — AKT 2: KLASYFIKACJA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Akt 2", color=ORANGE)
title_bar(s, "Odkrycie i klasyfikacja ryzyka")
bullets(s, 0.6, 2.0, 11.8, 3.2, [
    "Włączamy GovAI i rejestrujemy trzy agenty",
    "Asystent AI Act analizuje opis agenta kredytowego — klasyfikuje go automatycznie jako WYSOKIE RYZYKO (Aneks III, ocena zdolności kredytowej)",
    "Natychmiastowy raport: „Ten agent wymaga oceny zgodności przed kolejnym wdrożeniem”",
    "Konkretna lista działań z terminami — nie ogólne zalecenia, ale checklist",
], size=16, color=LGRAY, space_after=16)
footer(s, 13)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 14 — AKT 3: BRAMKA W DZIAŁANIU
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Akt 3", color=TEAL)
title_bar(s, "Bramka w działaniu — PII i blokady polityk")
rect(s, 0.6, 1.95, 5.85, 2.05, fill=BLUE, radius=True)
textbox(s, 0.85, 2.1, 5.4, 0.4, "ZAPYTANIE Z DANYMI OSOBOWYMI", size=11.5, color=TEAL, bold=True)
textbox(s, 0.85, 2.5, 5.4, 1.4,
        "„Oceń zdolność kredytową Jana Kowalskiego,\nPESEL 80010112345, ul. Mokotowska 15,\nWarszawa, zarobki 8500 zł netto.”",
        size=12.5, color=LGRAY, italic=True, line_spacing=1.2)
rect(s, 6.65, 1.95, 6.05, 2.05, fill=RGBColor(0x14, 0x2E, 0x24), radius=True)
textbox(s, 6.9, 2.1, 5.6, 0.4, "CO WIDZI MODEL PO MASKOWANIU", size=11.5, color=GREEN, bold=True)
textbox(s, 6.9, 2.5, 5.6, 1.4,
        "„Oceń zdolność kredytową [IMIĘ_NAZWISKO],\n[NUMER_ID], zamieszkałego w [ADRES]...”\n\n→ Presidio wykrył 3 kategorie PII",
        size=12.5, color=LGRAY, line_spacing=1.2)
rect(s, 0.6, 4.25, 12.1, 1.1, fill=RGBColor(0x3A, 0x1A, 0x1A), radius=True)
textbox(s, 0.9, 4.4, 11.5, 0.4, "PRÓBA NARUSZENIA POLITYKI", size=11.5, color=RED, bold=True)
textbox(s, 0.9, 4.75, 11.5, 0.5,
        "„Zmodyfikuj saldo konta klienta o -500 zł” → BRAMKA BLOKUJE W 87 ms",
        size=14, color=WHITE, bold=True)
textbox(s, 0.6, 5.55, 12, 0.6,
        "Dashboard: alert „Próba operacji poza zakresem agenta” + pełny wpis w dzienniku audytowym z metadanymi.",
        size=13, color=MGRAY, italic=True)
footer(s, 14)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 15 — AKT 4: NADZÓR CZŁOWIEKA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Akt 4", color=ORANGE)
title_bar(s, "Nadzór człowieka — art. 14 AI Act")
bullets(s, 0.6, 2.0, 7.0, 4.0, [
    "Agent kredytowy generuje decyzję odmowną — system wstrzymuje wykonanie",
    "Recenzent widzi: rekomendację agenta, uzasadnienie, dane źródłowe decyzji",
    "Trzy opcje: Zatwierdź / Zmień na pozytywną / Eskaluj",
    "Po kliknięciu „Zatwierdź” — dziennik zapisuje decyzję agenta, tożsamość recenzenta, czas przeglądu (28 sek.), decyzję końcową",
], size=15.5, space_after=16)
rect(s, 7.95, 2.0, 4.75, 3.9, fill=BLUE, radius=True)
textbox(s, 8.2, 2.25, 4.3, 0.8,
        "„Wasza firma może teraz powiedzieć regulatorowi:",
        size=14.5, color=TEAL, bold=True, italic=True, line_spacing=1.2)
textbox(s, 8.2, 3.0, 4.3, 1.6,
        "wiemy kto podjął każdą decyzję,\nkiedy i na podstawie jakich danych.”",
        size=16, color=WHITE, bold=True, italic=True, line_spacing=1.3)
footer(s, 15)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 16 — FINAŁ: RAPORT
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Finał", color=GREEN)
title_bar(s, "Raport dla regulatora — jeden klik")
rect(s, 0.6, 1.95, 12.1, 4.5, fill=RGBColor(0x10, 0x18, 0x2A), radius=True, line=TEAL)
textbox(s, 1.0, 2.2, 11, 0.5, "RAPORT OCENY ZGODNOŚCI Z USTAWĄ O SZTUCZNEJ INTELIGENCJI", size=14, color=TEAL, bold=True)
rows = [
    ("Agent 1 — Obsługa Klienta", "Ryzyko OGRANICZONE — zgodny ✓", GREEN),
    ("Agent 2 — Ocena Kredytowa", "Ryzyko WYSOKIE — wymaga 3 działań ⚠", ORANGE),
    ("Agent 3 — Rekrutacja", "Ryzyko WYSOKIE — wymaga 5 działań ⚠", ORANGE),
]
y = 2.85
for name, status, color in rows:
    textbox(s, 1.0, y, 5.5, 0.5, name, size=14, color=WHITE)
    textbox(s, 6.7, y, 5.5, 0.5, status, size=14, color=color, bold=True)
    y += 0.55
rect(s, 1.0, 4.65, 11.3, 0.02, fill=BLUE)
textbox(s, 1.0, 4.85, 11, 0.4, "PRIORYTETY NA NAJBLIŻSZE 30 DNI", size=13, color=TEAL, bold=True)
bullets(s, 1.0, 5.25, 11, 1.1, [
    "Ocena zgodności agenta kredytowego — termin: 14 dni",
    "Dokumentacja techniczna agenta rekrutacyjnego — termin: 21 dni",
], size=12.5, space_after=6)
footer(s, 16)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 17 — MODEL PARTNERSTWA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Pozycjonowanie")
title_bar(s, "Partnerstwo, nie zakup — firma doradcza jako kanał wdrożenia")
rect(s, 0.6, 2.0, 5.85, 3.7, fill=BLUE, radius=True)
textbox(s, 0.9, 2.2, 5.3, 0.4, "FIRMA KONSULTINGOWA / KANCELARIA", size=13, color=TEAL, bold=True)
bullets(s, 0.9, 2.7, 5.3, 2.8, [
    "Dostęp do klientów enterprise",
    "Wiedza o procesach regulacyjnych branży",
    "Zaufanie i relacje długoterminowe",
    "Zasoby do wdrożenia u klienta",
], size=13.5, space_after=12)
rect(s, 6.85, 2.0, 5.85, 3.7, fill=RGBColor(0x10, 0x2E, 0x2A), radius=True)
textbox(s, 7.15, 2.2, 5.3, 0.4, "GOVAI", size=13, color=TEAL, bold=True)
bullets(s, 7.15, 2.7, 5.3, 2.8, [
    "Technologia i produkt gotowy do wdrożenia",
    "Integracje i API niezależne od dostawcy modelu",
    "Szybkie wdrożenie — tygodnie, nie miesiące",
    "Aktualizacje prawne w miarę ewolucji AI Act",
], size=13.5, space_after=12)
textbox(s, 0.6, 5.95, 12, 1.2,
        "Model: firma doradcza sprzedaje i wdraża jako „powered by GovAI” (white-label lub branded) · "
        "podział przychodów negocjowany per kontrakt · wsparcie techniczne: GovAI, relacja z klientem: partner.",
        size=13, color=LGRAY, line_spacing=1.3)
footer(s, 17)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 18 — TRZY PYTANIA + RYZYKA
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
kicker(s, "Przed wyjściem z sali")
title_bar(s, "Trzy pytania, na które demo musi odpowiedzieć")
qa = [
    ("Czy to działa?", "Tak — żywe wywołania API, realnie wykryte dane osobowe, działający dziennik audytowy."),
    ("Czy można to wdrożyć u klienta?", "Tak — architektura kontenerowa, standardowe API, bez dostępu do kodu agentów klienta."),
    ("Czy jest pilność biznesowa?", "Tak — sierpień 2026, pierwsze kontrole AI Act, kary do 35 mln EUR. Sektor finansowy, prawny i medyczny już teraz powinien działać."),
]
y = 1.95
for q, a in qa:
    rect(s, 0.6, y, 12.1, 1.05, fill=BLUE, radius=True)
    textbox(s, 0.9, y + 0.1, 4.2, 0.8, q, size=14.5, color=TEAL, bold=True)
    textbox(s, 5.2, y + 0.12, 7.3, 0.85, a, size=12.5, color=LGRAY, line_spacing=1.15)
    y += 1.2
footer(s, 18)

# ════════════════════════════════════════════════════════════════════════════
# SLAJD 19 — NASTĘPNY KROK
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
textbox(s, 0.9, 1.6, 11.5, 0.5, "PROPONOWANY NASTĘPNY KROK", size=15, color=TEAL, bold=True)
rect(s, 0.9, 2.3, 11.5, 2.6, fill=BLUE, radius=True)
textbox(s, 1.25, 2.55, 10.8, 2.1,
        "„Dajcie nam jednego waszego klienta — instytucję finansową, kancelarię\n"
        "lub placówkę medyczną z minimum 20 agentami AI — na 6-tygodniowy pilotaż.\n"
        "My wdrażamy GovAI na ich środowisku testowym, wy asystujecie przy ocenie\n"
        "AI Act. Na koniec wspólnie prezentujemy wyniki ich zarządowi.”",
        size=17, color=WHITE, italic=True, bold=True, line_spacing=1.35)
textbox(s, 0.9, 5.25, 11.5, 0.9,
        "Pilotaż jest bezpłatny. Firma doradcza wychodzi z referencją i metodologią.\nGovAI wychodzi z referencją i umową produkcyjną.",
        size=15, color=LGRAY, line_spacing=1.3)
s.shapes.add_picture(LOGO_LOCKUP, Inches(0.9), Inches(6.55), height=Inches(0.45))
footer(s, 19)

prs.save(r"C:\Users\Adam\govai\presentations\GovAI_Prezentacja_Rynkowa.pptx")
print("OK — zapisano prezentację.")
